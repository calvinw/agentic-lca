#!/usr/bin/env python3
"""
Convert all .pptx files to Markdown and all .xlsx files to CSV
in the references/ directory, placing converted files alongside originals.
"""

import os
import csv
import glob

def convert_pptx_to_md(pptx_path):
    """Extract text from each slide and write a .md file alongside the original."""
    from pptx import Presentation
    from pptx.util import Pt

    md_path = os.path.splitext(pptx_path)[0] + ".md"
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"  SKIP (unreadable): {os.path.basename(pptx_path)} — {e}")
        return

    lines = [f"# {os.path.basename(os.path.splitext(pptx_path)[0])}\n"]

    for i, slide in enumerate(prs.slides, 1):
        slide_lines = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                # Detect heading vs body by font size if available
                level = 0
                try:
                    for run in para.runs:
                        if run.font.size and run.font.size >= Pt(20):
                            level = 2
                            break
                        elif run.font.size and run.font.size >= Pt(14):
                            level = 3
                            break
                except Exception:
                    pass
                if level == 2:
                    slide_lines.append(f"## {text}")
                elif level == 3:
                    slide_lines.append(f"### {text}")
                else:
                    slide_lines.append(f"- {text}")

        if slide_lines:
            lines.append(f"\n---\n\n<!-- Slide {i} -->\n")
            lines.extend(slide_lines)

    with open(md_path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines) + "\n")

    print(f"  OK: {os.path.basename(md_path)}")


def convert_xlsx_to_csv(xlsx_path):
    """Convert each sheet in an .xlsx file to its own .csv file alongside the original."""
    import openpyxl

    base = os.path.splitext(xlsx_path)[0]
    try:
        wb = openpyxl.load_workbook(xlsx_path, data_only=True)
    except Exception as e:
        print(f"  SKIP (unreadable): {os.path.basename(xlsx_path)} — {e}")
        return

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        # Sanitise sheet name for filename
        safe_name = "".join(c if c.isalnum() or c in " _-" else "_" for c in sheet_name).strip()
        csv_path = f"{base}__{safe_name}.csv"
        with open(csv_path, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            for row in ws.iter_rows(values_only=True):
                writer.writerow([("" if v is None else v) for v in row])
        print(f"  OK: {os.path.basename(csv_path)}")


def main():
    base_dir = os.path.dirname(os.path.abspath(__file__))

    pptx_files = glob.glob(os.path.join(base_dir, "**", "*.pptx"), recursive=True)
    xlsx_files = glob.glob(os.path.join(base_dir, "**", "*.xlsx"), recursive=True)

    print(f"\nFound {len(pptx_files)} .pptx files and {len(xlsx_files)} .xlsx files.\n")

    print("=== Converting PPTX → Markdown ===")
    for path in sorted(pptx_files):
        convert_pptx_to_md(path)

    print("\n=== Converting XLSX → CSV ===")
    for path in sorted(xlsx_files):
        convert_xlsx_to_csv(path)

    print("\nDone.")


if __name__ == "__main__":
    main()
